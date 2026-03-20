"""
Generate three PowerPoint presentations for the Visual ML Workflow Builder SIP project.

Run from ml_system/:
    python3 scripts/make_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0F, 0x34, 0x60)
DARK_NAVY  = RGBColor(0x16, 0x21, 0x3E)
CYAN       = RGBColor(0x00, 0xD4, 0xFF)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
ORANGE     = RGBColor(0xF3, 0x9C, 0x12)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE     = RGBColor(0x9B, 0x59, 0xB6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW     = RGBColor(0xF1, 0xC4, 0x0F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=DARK_NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=NAVY, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def label(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def title_slide(prs, title, subtitle, tag=None):
    sl = blank_slide(prs)
    bg(sl, DARK_NAVY)
    # Accent bar left
    box(sl, 0, 0, 0.5, 7.5, CYAN)
    # Title
    label(sl, title,   0.7, 1.8, 11.5, 1.6, size=40, bold=True,  color=CYAN,  align=PP_ALIGN.LEFT)
    label(sl, subtitle,0.7, 3.5, 11.0, 1.2, size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
    if tag:
        label(sl, tag, 0.7, 6.8, 6.0, 0.5, size=12, color=LIGHT_GRAY)
    return sl


def section_header(prs, title, color=CYAN):
    sl = blank_slide(prs)
    bg(sl, DARK_NAVY)
    box(sl, 0, 3.0, 13.33, 1.5, NAVY)
    label(sl, title, 0.5, 3.1, 12.0, 1.2, size=34, bold=True, color=color, align=PP_ALIGN.CENTER)
    return sl


def content_slide(prs, title, bullets, icon="", accent=CYAN, cols=1):
    sl = blank_slide(prs)
    bg(sl, DARK_NAVY)
    # Top bar
    box(sl, 0, 0, 13.33, 1.1, NAVY)
    label(sl, f"{icon}  {title}" if icon else title,
          0.3, 0.12, 12.5, 0.85, size=24, bold=True, color=accent)

    if cols == 1:
        y = 1.3
        for b in bullets:
            if b.startswith("##"):
                label(sl, b[2:].strip(), 0.5, y, 12.0, 0.45, size=15, bold=True, color=accent)
                y += 0.5
            elif b.startswith("--"):
                label(sl, "        " + b[2:].strip(), 0.5, y, 12.0, 0.4, size=13, color=LIGHT_GRAY)
                y += 0.42
            else:
                label(sl, "  •  " + b, 0.5, y, 12.0, 0.42, size=15, color=WHITE)
                y += 0.48
    else:
        mid = len(bullets) // 2
        for col_idx, col_bullets in enumerate([bullets[:mid], bullets[mid:]]):
            cx = 0.4 + col_idx * 6.5
            y = 1.3
            for b in col_bullets:
                if b.startswith("##"):
                    label(sl, b[2:].strip(), cx, y, 6.2, 0.42, size=14, bold=True, color=accent)
                    y += 0.48
                else:
                    label(sl, "  •  " + b, cx, y, 6.2, 0.42, size=13, color=WHITE)
                    y += 0.46
    return sl


def phase_slide(prs, phase_num, phase_name, status, items, status_color=GREEN):
    sl = blank_slide(prs)
    bg(sl, DARK_NAVY)
    box(sl, 0, 0, 13.33, 1.1, NAVY)
    label(sl, f"Phase {phase_num} — {phase_name}", 0.3, 0.12, 10.0, 0.85,
          size=24, bold=True, color=CYAN)
    label(sl, status, 10.5, 0.18, 2.5, 0.7, size=14, bold=True,
          color=status_color, align=PP_ALIGN.RIGHT)

    y = 1.35
    for item in items:
        if item.startswith("##"):
            label(sl, item[2:].strip(), 0.5, y, 12.0, 0.42, size=15, bold=True, color=CYAN)
            y += 0.52
        elif item.startswith("--"):
            label(sl, "          " + item[2:].strip(), 0.5, y, 12.0, 0.4, size=12, color=LIGHT_GRAY)
            y += 0.42
        else:
            chk = "✅" if "✅" in item else ("🔲" if "🔲" in item else "  •")
            text = item.replace("✅","").replace("🔲","").strip()
            label(sl, f"  {chk}  {text}", 0.5, y, 12.0, 0.42, size=14, color=WHITE)
            y += 0.48
    return sl


# ══════════════════════════════════════════════════════════════════════════════
#  DECK 1 — Annotation Tool
# ══════════════════════════════════════════════════════════════════════════════

def deck_annotation():
    prs = new_prs()

    title_slide(prs,
        "SVS Annotation Tool",
        "Visual ML Workflow Builder  |  Annotation Workflow Features",
        tag="UAT SIP Final Project 2026  |  Visual ML Workflow Builder")

    section_header(prs, "Drawing Modes")

    content_slide(prs, "Three Ways to Draw Annotations", [
        "## BBox Mode",
        "Click and drag to draw a bounding rectangle around an object",
        "Cursor: crosshair  |  Best for: fast rough labeling",
        "## Polygon Mode",
        "Click to place each vertex, double-click to close the shape",
        "Cursor: pencil  |  Best for: irregular shapes, precise outlines",
        "## SAM Box Mode",
        "Draw a rough box — SAM AI automatically tightens it to the object edge",
        "Cursor: tcross  |  Best for: complex outlines with one gesture",
    ], icon="✏️", accent=CYAN)

    section_header(prs, "SAM — Segment Anything Model")

    content_slide(prs, "AI-Assisted Segmentation Features", [
        "## Shift+Click → One-Shot SAM",
        "Click any point on an object — SAM generates a polygon automatically",
        "No drawing required — model finds the boundary",
        "## SAM Box Mode",
        "Rough box prompt → SAM tightens to a precise polygon in one step",
        "## ✦ Tighten with SAM",
        "Select any existing annotation → SAM re-segments it for a tighter fit",
        "Useful for correcting imprecise bounding boxes from YOLO pre-annotation",
    ], icon="🤖", accent=CYAN)

    content_slide(prs, "YOLO + SAM Pre-Annotation", [
        "## 🤖 Pre-annotate (YOLO bbox)",
        "Runs best.pt model on the current image",
        "Detected objects added as bounding boxes automatically",
        "Technician only needs to confirm, correct, or delete",
        "## 🤖✦ Pre-annotate (YOLO + SAM)",
        "YOLO detects objects + assigns class labels",
        "SAM immediately converts each detection to a tight polygon",
        "Result: full polygon annotations with correct labels — zero drawing",
        "Duplicate check: skips objects already annotated (>50% IoU overlap)",
    ], icon="🤖", accent=GREEN)

    section_header(prs, "Polygon Editing")

    content_slide(prs, "Edit Mode — Refine Any Polygon", [
        "Right-click a polygon → enters Edit Mode automatically",
        "Drag any vertex to reshape the polygon precisely",
        "Click on any edge → inserts a new vertex at that point",
        "Right-click a vertex → deletes that point",
        "Esc → exits edit mode and saves changes",
        "## Visual Cursors",
        "Fleur cursor  — idle in edit mode (ready to drag)",
        "Hand2 cursor  — hovering over a draggable vertex",
        "Pencil cursor — Polygon draw mode active",
    ], icon="🔷", accent=CYAN)

    section_header(prs, "Annotation Management")

    content_slide(prs, "Managing Labels", [
        "## Annotation List",
        "All annotations for the current image listed in the side panel",
        "Click any annotation to select and highlight it on the canvas",
        "## Relabel Selected",
        "Change the class of any existing annotation without redrawing",
        "## Delete Selected",
        "Remove incorrect annotations instantly (also: Delete/Backspace keys)",
        "## Active Class Lock",
        "Selected drawing class stays locked — clicking annotation list won't reset it",
        "Prevents accidental mislabeling of new annotations",
    ], icon="🏷️", accent=ORANGE)

    section_header(prs, "Navigation & Productivity")

    content_slide(prs, "Image Browser & Workflow Tools", [
        "## Image Browser Panel",
        "Scrollable list of all images in the project on the right panel",
        "Search bar — filter by filename instantly",
        "Filter buttons: All  |  ✓ Done  |  ○ Todo  |  ⚑ Flagged",
        "## ⏭ Skip to Unannotated",
        "Jumps directly to the next image with no annotations — saves time",
        "## ⚑ Flag Image",
        "Mark an image for later review (bad photo, uncertain label, needs expert)",
        "## H — Toggle Annotations",
        "Hide/show all annotation overlays — compare clean vs labeled view",
    ], icon="🗂️", accent=CYAN)

    content_slide(prs, "Visual Aids", [
        "## Brightness Slider",
        "Adjust image brightness on-the-fly for shadowed or dark photos",
        "SAM uses the brightness-adjusted image for better edge detection",
        "Reset button returns to original brightness",
        "## 📊 Class Balance Chart",
        "Live bar chart showing annotation count per class across the dataset",
        "Instantly reveals underrepresented classes that need more labeling",
        "## Confidence Score Stored",
        "Pre-annotated items store the model confidence score in the JSON",
        "Enables filtering by confidence for review prioritisation",
    ], icon="🎨", accent=PURPLE)

    section_header(prs, "Export & Reporting")

    content_slide(prs, "Export Features", [
        "## 📷 Export Annotated Image",
        "Renders current image at full resolution with annotation overlays",
        "Thick outlines with dark shadow pass for visibility on any background",
        "Semi-transparent fills show object regions clearly",
        "Labels rendered at scaled font size — readable on high-res photos",
        "Saved to: projects/svs_plumbing/reports/annotated/",
        "## Auto-Save",
        "All changes saved to COCO JSON automatically on every action",
        "No manual save required — no work is ever lost",
    ], icon="💾", accent=YELLOW)

    return prs


# ══════════════════════════════════════════════════════════════════════════════
#  DECK 2 — Full Technician Workflow
# ══════════════════════════════════════════════════════════════════════════════

def deck_technician():
    prs = new_prs()

    title_slide(prs,
        "SVS Technician Workflow",
        "AI-Assisted Compliance Inspection  |  Active Learning Pipeline",
        tag="UAT SIP Final Project 2026  |  Visual ML Workflow Builder")

    section_header(prs, "The Problem We're Solving", color=ORANGE)

    content_slide(prs, "Current SVS Inspection Challenges", [
        "Technicians photograph 21 steps of every ice machine installation",
        "Photos reviewed manually — slow, inconsistent, human error",
        "No automated check that all required components are present",
        "No way to verify installation_date, serial numbers, filter specs",
        "Model accuracy limited by available training data — improves slowly",
    ], icon="⚠️", accent=ORANGE)

    section_header(prs, "The Solution", color=GREEN)

    content_slide(prs, "AI-Assisted Inspection Loop", [
        "## Step 1 — Technician Takes Photos",
        "21 photos per installation using phone or tablet",
        "Uploaded directly through the Laravel web app",
        "## Step 2 — Model Inspects Automatically",
        "FastAPI microservice runs YOLO detection on each photo",
        "Returns: detected components, confidence scores, missing items",
        "## Step 3 — Technician Reviews Results",
        "Web UI shows each photo with bounding box overlays",
        "Green = confirmed  |  Yellow = low confidence  |  Red = missing",
    ], icon="🔄", accent=GREEN)

    content_slide(prs, "Active Learning — Model Improves Over Time", [
        "## Step 4 — Technician Corrects Mistakes",
        "Confirms correct detections with one tap",
        "Draws box around anything the model missed",
        "Relabels any wrong predictions",
        "## Step 5 — Corrections Queued",
        "Every correction stored in the training queue automatically",
        "No extra work for the technician — feedback happens during normal use",
        "## Step 6 — Automatic Retraining",
        "When queue reaches threshold → new training run triggered",
        "New best.pt deployed → model immediately more accurate",
    ], icon="🧠", accent=CYAN)

    content_slide(prs, "The Active Learning Loop", [
        "## The Flywheel Effect",
        "More technicians using the app → more corrections → better model",
        "Better model → fewer corrections needed → faster inspections",
        "Faster inspections → more adoption → even more data",
        "## What the Technician Experiences",
        "Day 1: model misses ~25% of objects, tech corrects them",
        "Month 1: model misses ~15%, corrections take 30 seconds",
        "Month 3: model misses ~5%, mostly just confirm and submit",
        "## What You Get",
        "A model trained on real SVS equipment by domain experts",
        "Continuously improving without manual annotation sessions",
    ], icon="🔁", accent=GREEN)

    content_slide(prs, "Inspection Output — Pass / Fail / Report", [
        "## Detection Layer (YOLO)",
        "Locates: water_filter, pressure_gauge, water_line, water_valve, drain_line, disconnect_box",
        "Returns bounding box + class + confidence for each detected component",
        "## Reading Layer (OCR)",
        "Crops detected text regions: installation_date, data_sticker, qr_code",
        "Reads actual values: dates, serial numbers, filter specifications",
        "## Inspection Engine",
        "Checks all required components present  ✓ / ✗",
        "Validates dates, serial numbers against compliance rules",
        "Generates score + Pass/Fail + PDF report → sent to Laravel",
    ], icon="📋", accent=ORANGE)

    content_slide(prs, "System Components", [
        "## Frontend — Laravel Web App",
        "Technician login, photo upload, results review, report download",
        "## Microservice — FastAPI",
        "POST /inspect   → run model, return detections",
        "POST /correct   → store technician corrections to queue",
        "GET  /status    → model version, accuracy, queue size",
        "## ML Training Suite — ml_system/",
        "annotate.py, train.py, review.py, convert_to_yolo.py",
        "## Compute — Home GPU or RunPod",
        "RTX 4070 Super (home) or RunPod API for cloud GPU training",
    ], icon="🏗️", accent=CYAN, cols=1)

    return prs


# ══════════════════════════════════════════════════════════════════════════════
#  DECK 3 — Build Order & Phases
# ══════════════════════════════════════════════════════════════════════════════

def deck_build_order():
    prs = new_prs()

    title_slide(prs,
        "Build Order & Phases",
        "Visual ML Workflow Builder  |  Development Roadmap",
        tag="UAT SIP Final Project 2026  |  Visual ML Workflow Builder")

    section_header(prs, "Overview — 6 Phases", color=CYAN)

    content_slide(prs, "Roadmap at a Glance", [
        "## ✅ Phase 1 — Local ML Training Suite         COMPLETE",
        "annotation tool, training pipeline, model review",
        "## 🔲 Phase 2 — FastAPI Inference Microservice",
        "wrap best.pt in an API, serve detections over HTTP",
        "## 🔲 Phase 3 — Laravel Integration",
        "connect Laravel to FastAPI, display detections in UI",
        "## 🔲 Phase 4 — Correction Loop & Retraining",
        "technician corrections queue → automated retraining trigger",
        "## 🔲 Phase 5 — OCR + Inspection Engine",
        "read text from labels, generate compliance score + report",
        "## 🔲 Phase 6 — Cloud GPU & Scale",
        "RunPod API, distributed training, multi-project support",
    ], icon="🗺️", accent=CYAN)

    phase_slide(prs, 1, "Local ML Training Suite", "✅ COMPLETE", [
        "✅ annotate.py — full annotation tool with SAM, brightness, class balance",
        "✅ Drawing modes: BBox, Polygon, SAM Box",
        "✅ Edit mode: drag vertices, insert/delete points",
        "✅ Pre-annotation: YOLO bbox and YOLO+SAM polygon",
        "✅ Image browser with search + filter (All/Done/Todo/Flagged)",
        "✅ Export annotated image at full resolution",
        "✅ convert_to_yolo.py — COCO JSON → YOLO segmentation format",
        "✅ train.py — experiment runner, saves config + results JSON",
        "✅ review.py — model performance review, surfaces training plots",
        "✅ utils.py — mAP comparison chart across experiments",
        "✅ exp_001 trained — mAP@50: 0.75  |  Precision: 0.84  |  Recall: 0.71",
    ], status_color=GREEN)

    phase_slide(prs, 2, "FastAPI Inference Microservice", "🔲 NEXT", [
        "## Goal: bridge ml_system to Laravel over HTTP",
        "🔲 scripts/api.py — FastAPI app",
        "🔲 POST /inspect — accepts base64 image, returns detections JSON",
        "--  { class, confidence, bbox } for each detected object",
        "--  missing components list",
        "--  overall score (% of required components found)",
        "🔲 GET /status — model version, mAP, queue size",
        "🔲 Model hot-reload when new best.pt is deployed",
        "🔲 Docker container for easy deployment",
        "## Estimated effort: 2–3 days",
    ], status_color=YELLOW)

    phase_slide(prs, 3, "Laravel Integration", "🔲 PLANNED", [
        "## Goal: technicians see detection results in the web app",
        "🔲 InspectionController — sends photo to FastAPI, stores result",
        "🔲 Detection overlay UI — bounding boxes drawn on photo in browser",
        "--  Green box = high confidence  |  Yellow = low  |  Red = missing",
        "🔲 Confirm / Correct UI — technician reviews each detection",
        "--  Tap to confirm  |  Draw to correct  |  Add missed objects",
        "🔲 POST /correct — corrections sent back to FastAPI queue",
        "🔲 Admin dashboard — queue size, model version, accuracy",
        "## Estimated effort: 1 week",
    ], status_color=YELLOW)

    phase_slide(prs, 4, "Correction Loop & Automated Retraining", "🔲 PLANNED", [
        "## Goal: technician corrections automatically improve the model",
        "🔲 Correction queue stored in COCO JSON format",
        "🔲 Threshold trigger — retrain when N corrections accumulated",
        "🔲 convert_to_yolo.py runs automatically on new dataset version",
        "🔲 train.py dispatched to home GPU server or RunPod API",
        "🔲 New best.pt pulled back and deployed to FastAPI",
        "🔲 Model version tracked — rollback if accuracy drops",
        "🔲 Notification when retraining completes",
        "## Estimated effort: 1–2 weeks",
    ], status_color=LIGHT_GRAY)

    phase_slide(prs, 5, "OCR + Inspection Engine", "🔲 PLANNED", [
        "## Goal: read text values and generate compliance score",
        "🔲 OCR integration (easyocr or paddleocr)",
        "--  Crops YOLO-detected text regions: installation_date, data_sticker, qr_code",
        "--  Returns actual text values from the image",
        "🔲 Inspection Engine",
        "--  Checks all required components present",
        "--  Validates dates, serial numbers against compliance rules",
        "--  Generates score: X/21 steps passed",
        "🔲 Report generation — PDF with photo evidence + scores",
        "🔲 Pass/Fail result returned to Laravel",
        "## Estimated effort: 1–2 weeks",
    ], status_color=LIGHT_GRAY)

    phase_slide(prs, 6, "Cloud GPU & Scale", "🔲 FUTURE", [
        "## Goal: fast retraining, multi-project, production-ready",
        "🔲 Home GPU server — Flask API on RTX 4070 Super",
        "--  POST /train, GET /status, GET /weights",
        "--  Tailscale for secure remote access",
        "🔲 RunPod API integration — cloud GPU fallback",
        "--  Automatic dataset upload, training job submission",
        "--  Cost tracking per experiment",
        "🔲 Multi-project support — beyond SVS plumbing",
        "🔲 Electron app shell wrapping the full workflow",
        "🔲 Team features — multiple annotators, role-based access",
        "## Estimated effort: 2–4 weeks",
    ], status_color=LIGHT_GRAY)

    content_slide(prs, "What's Already Built vs What's Next", [
        "## Built (Phase 1)",
        "Full annotation pipeline from raw photos to trained YOLO model",
        "SAM AI-assisted polygon annotation with YOLO pre-annotation",
        "Model review and experiment comparison tooling",
        "First trained model: mAP@50 = 0.75 on 6 SVS component classes",
        "## Immediate Next Step (Phase 2)",
        "FastAPI endpoint — 2-3 days of focused work",
        "Single file: scripts/api.py — bridges everything to Laravel",
        "## The Leverage Point",
        "Once Phase 2 + 3 are done, every technician inspection becomes training data",
        "Model improves automatically — no manual annotation sessions needed",
    ], icon="📍", accent=GREEN)

    return prs


# ── Generate all three ────────────────────────────────────────────────────────

OUTPUT_DIR = os.path.join("projects", "svs_plumbing", "reports", "slides")
os.makedirs(OUTPUT_DIR, exist_ok=True)

def deck_dataset_expansion():
    prs = new_prs()

    title_slide(prs,
        "Dataset Expansion Strategy",
        "Semi-Supervised Annotation  |  Batch Pre-Annotation Workflow",
        tag="UAT SIP Final Project 2026  |  Visual ML Workflow Builder")

    section_header(prs, "The Core Idea", color=CYAN)

    content_slide(prs, "Semi-Supervised Annotation", [
        "You cannot skip manual annotation entirely — the model must learn first",
        "But you can dramatically reduce how much manual work is needed",
        "## The 10x Rule",
        "Manually annotate 10–15 examples of a new class",
        "Retrain — model now recognises the class",
        "Let YOLO+SAM pre-annotate 100+ photos automatically",
        "You only review and correct — not draw from scratch",
        "## Why This Works",
        "YOLO generalises from a small number of good examples",
        "SAM handles the precise polygon shape automatically",
        "Your time shifts from drawing → reviewing (10x faster)",
    ], icon="🧠", accent=CYAN)

    content_slide(prs, "Why You Must Retrain First", [
        "## YOLO can only find what it was trained to recognise",
        "A new class does not exist in the model until it has seen examples",
        "Running batch pre-annotate before retraining = zero detections",
        "## The Introduction Step",
        "10–15 manual annotations = introducing the class to the model",
        "Each example shows YOLO: shape, colour, context, typical location",
        "After retraining, YOLO can generalise to similar objects it has never seen",
        "## Minimum Examples to Retrain",
        "10–15 annotations → model will detect the class but may miss some",
        "30–50 annotations → reliable detection, ready for batch pre-annotation",
        "100+ annotations → strong model, high confidence on new images",
    ], icon="⚠️", accent=ORANGE)

    section_header(prs, "The Expansion Cycle", color=GREEN)

    content_slide(prs, "Step-by-Step Dataset Expansion", [
        "## Step 1 — Manually Annotate 10–15 New Examples",
        "Use Polygon or SAM Box mode for accurate shapes",
        "Vary the examples: different lighting, angles, distances",
        "Include edge cases — partial views, shadows, overlapping objects",
        "## Step 2 — Convert Dataset + Retrain",
        "python3 projects/svs_plumbing/convert_to_yolo.py --version v2",
        "python3 scripts/train.py --project svs_plumbing",
        "Takes ~3 hrs CPU  or  ~15 min on RTX 4070 / Colab GPU",
        "## Step 3 — Review Model Performance",
        "python3 scripts/review.py --project svs_plumbing",
        "Check mAP on new class — if above 0.40, proceed to batch",
    ], icon="1️⃣", accent=GREEN)

    content_slide(prs, "Step-by-Step Dataset Expansion (cont.)", [
        "## Step 4 — Batch Pre-Annotate 100+ Photos",
        "Click '🤖✦ Batch Pre-annotate (All Unannotated)'",
        "Processes every unannotated image in the dataset automatically",
        "Status bar shows live progress: 'Processing 47/312...'",
        "Takes a few minutes — runs in background",
        "## Step 5 — Review Pass",
        "Filter image browser: ○ Todo → shows newly pre-annotated images",
        "Browse quickly: confirm good boxes, delete wrong ones, add missed objects",
        "Target: spend < 10 seconds per image on review",
        "## Step 6 — Retrain Again",
        "Convert + train on the full expanded dataset",
        "Model now has hundreds of examples — much stronger accuracy",
    ], icon="2️⃣", accent=GREEN)

    content_slide(prs, "The Compounding Effect", [
        "## Cycle 1 — Introducing a new class",
        "15 manual annotations → retrain → batch 100 → review → retrain",
        "End result: ~115 annotations, model mAP ~0.55–0.65",
        "## Cycle 2 — Strengthening",
        "Batch 200 more → quick review → retrain",
        "End result: ~315 annotations, model mAP ~0.70–0.80",
        "## Cycle 3 — Production quality",
        "Batch remaining images → minimal review → retrain",
        "End result: 500+ annotations, model mAP ~0.80–0.90",
        "## Time Comparison",
        "Manual only:  500 annotations × 2 min each = 16+ hours",
        "This method:  15 manual + 3 review passes  = 2–3 hours total",
    ], icon="📈", accent=CYAN)

    content_slide(prs, "Adding Multiple New Classes", [
        "## One Class at a Time vs All at Once",
        "Adding one class per cycle is safer — easier to diagnose problems",
        "Adding multiple classes works if you have 10–15 examples of each",
        "## Recommended Order for SVS Dataset",
        "Priority 1: classes with lowest mAP in review.py output",
        "Priority 2: classes most critical for compliance checking",
        "Priority 3: classes that appear most frequently in new photos",
        "## Current SVS Classes (6 trained)",
        "water_line, water_valve, drain_line, filter, pressure_gauge, disconnect_box",
        "## Remaining 19 Classes to Expand Into",
        "data_sticker, qr_code, installation_date, floor_drain, air-gap_fitting,",
        "water_filter, pressure_gauge, bin_drain_line, BackOfMachine, IceBin, ...",
    ], icon="🗂️", accent=PURPLE)

    content_slide(prs, "Batch Pre-Annotation in the Tool", [
        "## New Button: 🤖✦ Batch Pre-annotate (All Unannotated)",
        "Finds every image in the dataset with no annotations",
        "Runs YOLO detection on each image",
        "Feeds each YOLO bbox into SAM for tight polygon",
        "Saves polygon annotations with class labels to COCO JSON",
        "Skips images that already have annotations (safe to re-run)",
        "## Progress Feedback",
        "Status bar: 'Batch: 47/312 done — 12 added, 2 skipped'",
        "Can be stopped at any time — work already done is saved",
        "## After Batch Completes",
        "Image browser filter ○ Todo shows what still needs review",
        "Use ✓ Done filter to see what was successfully pre-annotated",
    ], icon="⚡", accent=GREEN)

    return prs


def deck_exact_flow():
    prs = new_prs()

    title_slide(prs,
        "The Exact Training Cycle",
        "Annotate  →  Export  →  Colab GPU  →  Deploy  →  Batch  →  Repeat",
        tag="UAT SIP Final Project 2026  |  Visual ML Workflow Builder")

    # Overview
    content_slide(prs, "The Full Cycle at a Glance", [
        "## Step 1 — Annotate New Photos Locally",
        "Add new images, label 10–15 examples of each new class in annotate.py",
        "## Step 2 — Export Dataset",
        "python3 projects/svs_plumbing/convert_to_yolo.py --version v2",
        "## Step 3 — Train on Google Colab GPU",
        "Zip dataset → upload to Google Drive → run Colab notebook",
        "Free T4 GPU — trains in 30–45 min instead of 3 hours on CPU",
        "## Step 4 — Deploy New Model",
        "Download best.pt from Colab → drop into ml_system/models/best.pt",
        "## Step 5 — Batch Pre-annotate + Review",
        "⚡ Batch Pre-annotate All → quick review pass → convert v3 → repeat",
    ], icon="🔄", accent=CYAN)

    # Step 1
    content_slide(prs, "Step 1 — Annotate New Photos Locally", [
        "## Add New Images to the Dataset",
        "Copy new photos into:",
        "-- projects/svs_plumbing/datasets/raw/",
        "## Open the Annotation Tool",
        "-- DISPLAY=:0 python3 scripts/annotate.py --project svs_plumbing",
        "## Label 10–15 Examples of Each New Class",
        "Use SAM Box mode (tcross cursor) for fast accurate polygons",
        "Use Shift+click for one-shot SAM on clear objects",
        "Use Polygon mode for complex shapes",
        "## Tips for Good Examples",
        "Vary lighting, angles, distances across your 10–15 examples",
        "Include edge cases: partial views, shadows, overlapping objects",
        "Use the Brightness slider for dark/shadowed images",
    ], icon="1️⃣", accent=CYAN)

    # Step 2
    content_slide(prs, "Step 2 — Export Dataset", [
        "## Run the Converter",
        "From ml_system/ run:",
        "-- python3 projects/svs_plumbing/convert_to_yolo.py --version v2",
        "## What It Does",
        "Reads COCO JSON → converts to YOLO segmentation format",
        "Splits images into train (85%) and val (15%) sets",
        "Writes data.yaml with absolute paths YOLO can read",
        "## Output Location",
        "-- projects/svs_plumbing/datasets/dataset_v2/",
        "--   images/train/   images/val/",
        "--   labels/train/   labels/val/",
        "--   data.yaml",
        "## Version Each Export",
        "Use --version v2, v3, v4 to keep history of each training cycle",
    ], icon="2️⃣", accent=ORANGE)

    # Step 3
    content_slide(prs, "Step 3 — Train on Google Colab GPU", [
        "## Why Colab Instead of CPU",
        "CPU training (Intel Core Ultra 7):  ~3 hours per 50 epochs",
        "Colab T4 GPU (free):               ~30–45 min per 50 epochs",
        "Home RTX 4070 Super:               ~12–15 min per 50 epochs",
        "## Colab Workflow",
        "1. Zip dataset_v2/ folder on your machine",
        "2. Upload zip to Google Drive",
        "3. Open Colab notebook",
        "4. Mount Google Drive → unzip dataset",
        "5. Run: model.train(data='data.yaml', epochs=50, imgsz=640)",
        "6. Download best.pt from runs/train/weights/best.pt",
        "## Colab Tips",
        "Use GPU runtime: Runtime → Change runtime type → T4 GPU",
        "Keep browser tab active — Colab disconnects if idle",
    ], icon="3️⃣", accent=GREEN)

    # Step 4
    content_slide(prs, "Step 4 — Deploy New Model", [
        "## The One File to Remember",
        "Drop your downloaded best.pt into exactly this path:",
        "-- ml_system/models/best.pt",
        "## Why This Location",
        "It is the first path _find_best_pt() checks",
        "Every tool reads from here automatically — one swap upgrades everything",
        "## What Gets Upgraded Instantly",
        "🤖  Pre-annotate (YOLO bbox)      — uses new model",
        "🤖✦ Pre-annotate (YOLO+SAM)       — uses new model",
        "⚡  Batch Pre-annotate All        — uses new model",
        "review.py                          — reviews new model",
        "## No Code Changes Needed",
        "Just drop the file — the system picks it up on next button click",
    ], icon="4️⃣", accent=YELLOW)

    # Step 5
    content_slide(prs, "Step 5 — Batch Pre-annotate + Review", [
        "## Run Batch Pre-annotation",
        "In annotate.py click: ⚡ Batch Pre-annotate All",
        "Finds every image with no annotations",
        "Runs YOLO detection + SAM polygon on each image automatically",
        "Status bar shows live progress: 'Batch: 47/312 — 180 polygons added'",
        "Saves every 10 images — safe to stop and resume",
        "## Do a Review Pass",
        "Filter image browser: ○ Todo → shows newly pre-annotated images",
        "Browse quickly — target under 10 seconds per image",
        "Confirm good boxes, delete wrong ones, draw any missed objects",
        "## Then Repeat the Cycle",
        "convert_to_yolo.py --version v3 → Colab → best.pt → batch → review",
        "Each cycle: model gets stronger, review pass gets faster",
    ], icon="5️⃣", accent=CYAN)

    # The repeat loop
    content_slide(prs, "The Repeating Cycle", [
        "## Cycle 1 — Introducing New Classes",
        "15 manual annotations → export v2 → Colab train → deploy → batch 100+",
        "Review pass: ~30 min  |  Result: model knows new classes",
        "## Cycle 2 — Strengthening",
        "Batch 200 more → review → export v3 → Colab → deploy",
        "Review pass: ~15 min  |  Result: mAP climbs toward 0.80+",
        "## Cycle 3 — Production Quality",
        "Batch remaining images → minimal review → export v4 → final train",
        "Review pass: ~5 min   |  Result: 500+ examples, ready for deployment",
        "## Time Investment Per Cycle",
        "Manual annotation:    30–60 min (10–15 new examples)",
        "Colab training:       30–45 min (unattended)",
        "Batch pre-annotate:   10–20 min (unattended)",
        "Review pass:          15–30 min (shrinks each cycle)",
    ], icon="🔁", accent=GREEN)

    # The single file rule
    content_slide(prs, "The Golden Rule — One File", [
        "## ml_system/models/best.pt",
        "This single file is the brain of the entire system",
        "Every annotation button, every batch run, every review reads from here",
        "Upgrading the model = dropping one file",
        "## Version Your Models",
        "Keep copies with meaningful names before overwriting:",
        "-- best_exp001_mAP75.pt",
        "-- best_exp002_mAP81.pt",
        "-- best_exp003_mAP87.pt",
        "Then copy the best one to models/best.pt",
        "## If a New Model Performs Worse",
        "Just copy the previous version back to models/best.pt",
        "Instant rollback — no code changes needed",
    ], icon="🔑", accent=YELLOW)

    return prs


if __name__ == "__main__":
    decks = [
        ("01_Annotation_Tool.pptx",         deck_annotation),
        ("02_Technician_Workflow.pptx",     deck_technician),
        ("03_Build_Order_Phases.pptx",      deck_build_order),
        ("04_Dataset_Expansion.pptx",       deck_dataset_expansion),
        ("05_Exact_Training_Cycle.pptx",    deck_exact_flow),
    ]

    for filename, builder in decks:
        path = os.path.join(OUTPUT_DIR, filename)
        try:
            prs = builder()
            prs.save(path)
            print(f"Saved: {path}")
        except PermissionError:
            print(f"SKIPPED (file open in PowerPoint): {filename}")

    print(f"\nDone. Slides at: {os.path.abspath(OUTPUT_DIR)}")
